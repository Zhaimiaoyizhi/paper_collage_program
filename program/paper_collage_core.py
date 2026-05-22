from __future__ import annotations

import argparse
import csv
import io
import json
import sys
import tarfile
import time
import urllib.request
import xml.etree.ElementTree as ET
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Iterable
from urllib.parse import quote, urljoin, urlparse

import fitz
import requests
from bs4 import BeautifulSoup
from PIL import Image, ImageColor, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor


USER_AGENT = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
    "AppleWebKit/537.36 (KHTML, like Gecko) "
    "Chrome/124.0 Safari/537.36"
)


@dataclass(frozen=True)
class PaperSource:
    slug: str
    refs: str
    article_url: str
    pdf_url: str | None = None


@dataclass(frozen=True)
class DrawPlacement:
    source_index: int
    group_index: int
    layer_index: int
    paste_xy: tuple[int, int]


@dataclass(frozen=True)
class RenderedPage:
    source: PaperSource
    placement: DrawPlacement
    framed_path: Path
    paste_xy: tuple[int, int]
    size: tuple[int, int]


@dataclass(frozen=True)
class InteractiveLoginConfig:
    enabled: bool = False
    profile_dir: Path = Path(".playwright_login_profile")
    timeout_seconds: int = 900
    poll_interval_seconds: int = 10


class MissingPdfError(RuntimeError):
    def __init__(self, missing: list[PaperSource]):
        self.missing = missing
        super().__init__(f"{len(missing)} PDFs are missing")


PDF_SOURCES: list[PaperSource] = [
    PaperSource(
        "p01_s41586_024_07473_2",
        "[1] [2] [4] [6] [12] [24]",
        "https://www.nature.com/articles/s41586-024-07473-2",
        "https://www.nature.com/articles/s41586-024-07473-2.pdf",
    ),
    PaperSource(
        "p02_s41586_023_06457_y",
        "[3] [5] [13] [18]",
        "https://www.nature.com/articles/s41586-023-06457-y",
        "https://www.nature.com/articles/s41586-023-06457-y.pdf",
    ),
    PaperSource(
        "p03_nature24018",
        "[7] [14] [15]",
        "https://www.nature.com/articles/nature24018",
        "https://www.nature.com/articles/nature24018.pdf",
    ),
    PaperSource(
        "p04_nature13206",
        "[8] [25] [26]",
        "https://www.nature.com/articles/nature13206",
        "https://www.nature.com/articles/nature13206.pdf",
    ),
    PaperSource(
        "p05_s10577_011_9252_1",
        "[9]",
        "https://link.springer.com/article/10.1007/s10577-011-9252-1",
        "https://link.springer.com/content/pdf/10.1007/s10577-011-9252-1.pdf",
    ),
    PaperSource(
        "p06_pmc9894122",
        "[10] [23]",
        "https://pmc.ncbi.nlm.nih.gov/articles/PMC9894122/",
        "https://pmc.ncbi.nlm.nih.gov/articles/PMC9894122/pdf/",
    ),
    PaperSource(
        "p07_s41559_022_01974_x",
        "[11] [17]",
        "https://www.nature.com/articles/s41559-022-01974-x",
        "https://www.nature.com/articles/s41559-022-01974-x.pdf",
    ),
    PaperSource(
        "p08_genome_30_6_860",
        "[16]",
        "https://genome.cshlp.org/content/30/6/860.full.pdf",
        "https://genome.cshlp.org/content/30/6/860.full.pdf",
    ),
    PaperSource(
        "p09_s41586_023_06425_6",
        "[19]",
        "https://www.nature.com/articles/s41586-023-06425-6",
        "https://www.nature.com/articles/s41586-023-06425-6.pdf",
    ),
    PaperSource(
        "p10_s41576_024_00757_3",
        "[20]",
        "https://www.nature.com/articles/s41576-024-00757-3",
        "https://www.nature.com/articles/s41576-024-00757-3.pdf",
    ),
    PaperSource(
        "p11_s41594_024_01362_y",
        "[21]",
        "https://www.nature.com/articles/s41594-024-01362-y",
        "https://www.nature.com/articles/s41594-024-01362-y.pdf",
    ),
    PaperSource(
        "p12_s41559_024_02627_x",
        "[22]",
        "https://www.nature.com/articles/s41559-024-02627-x",
        "https://www.nature.com/articles/s41559-024-02627-x.pdf",
    ),
]


def safe_name_from_url(url: str) -> str:
    parsed = urlparse(url)
    raw = f"{parsed.netloc}{parsed.path}".strip("/")
    chars: list[str] = []
    for char in raw.lower():
        chars.append(char if char.isalnum() else "_")
    collapsed = "_".join(part for part in "".join(chars).split("_") if part)
    return collapsed


def normalize_slug(value: str) -> str:
    chars: list[str] = []
    for char in value.strip().lower():
        chars.append(char if char.isalnum() else "_")
    collapsed = "_".join(part for part in "".join(chars).split("_") if part)
    return collapsed or "paper"


def pick_first_nonempty(row: dict[str, str], field_names: list[str]) -> str:
    for field_name in field_names:
        value = (row.get(field_name) or "").strip()
        if value:
            return value
    return ""


def load_sources_from_csv(csv_path: Path) -> list[PaperSource]:
    with csv_path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        if reader.fieldnames is None:
            raise ValueError(f"{csv_path} has no header row")

        sources: list[PaperSource] = []
        seen_urls: set[str] = set()
        for row_index, raw_row in enumerate(reader, start=1):
            row = {str(key).strip(): (value or "") for key, value in raw_row.items() if key is not None}
            article_url = pick_first_nonempty(row, ["article_url", "url", "paper_url", "website", "article"])
            if not article_url:
                raise ValueError(
                    f"Row {row_index} in {csv_path} is missing article_url/url/paper_url/website"
                )
            if article_url in seen_urls:
                continue
            seen_urls.add(article_url)

            raw_slug = pick_first_nonempty(row, ["slug", "id", "name"])
            slug = normalize_slug(raw_slug) if raw_slug else f"p{len(sources) + 1:02d}_{safe_name_from_url(article_url)}"
            refs = pick_first_nonempty(row, ["refs", "ref", "reference", "label"]) or f"[{row_index}]"
            pdf_url = pick_first_nonempty(row, ["pdf_url", "direct_pdf_url", "download_url"]) or None
            sources.append(PaperSource(slug=slug, refs=refs, article_url=article_url, pdf_url=pdf_url))

    if not sources:
        raise ValueError(f"{csv_path} contains no usable paper rows")
    return sources


def load_sources_from_txt(txt_path: Path) -> list[PaperSource]:
    sources: list[PaperSource] = []
    seen_urls: set[str] = set()
    for raw_line in txt_path.read_text(encoding="utf-8-sig").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#"):
            continue
        article_url = line
        if article_url in seen_urls:
            continue
        seen_urls.add(article_url)
        slug = f"p{len(sources) + 1:02d}_{safe_name_from_url(article_url)}"
        sources.append(
            PaperSource(slug=slug, refs=f"[{len(sources) + 1}]", article_url=article_url, pdf_url=None)
        )

    if not sources:
        raise ValueError(f"{txt_path} contains no usable paper URLs")
    return sources


def load_sources_from_input_file(input_path: Path) -> list[PaperSource]:
    suffix = input_path.suffix.lower()
    if suffix == ".csv":
        return load_sources_from_csv(input_path)
    if suffix == ".txt":
        return load_sources_from_txt(input_path)
    raise ValueError(f"Unsupported input format for {input_path}. Expected .csv or .txt")


def infer_pdf_candidates(article_url: str) -> list[str]:
    parsed = urlparse(article_url)
    url = article_url.rstrip("/")
    candidates: list[str] = []

    if article_url.lower().endswith(".pdf"):
        candidates.append(article_url)
    elif parsed.netloc.endswith("nature.com") and "/articles/" in parsed.path:
        candidates.append(f"{url}.pdf")
    elif parsed.netloc.endswith("springer.com") and "/article/" in parsed.path:
        doi = parsed.path.split("/article/", 1)[1].strip("/")
        candidates.append(f"https://link.springer.com/content/pdf/{doi}.pdf")
    elif parsed.netloc.endswith("ncbi.nlm.nih.gov") and "/articles/" in parsed.path:
        candidates.append(f"{url}/pdf/")

    candidates.append(article_url)
    return unique_ordered(candidates)


def infer_landing_page_candidates(article_url: str) -> list[str]:
    parsed = urlparse(article_url)
    lowered_path = parsed.path.lower()
    candidates: list[str] = []

    if lowered_path.endswith(".full.pdf"):
        landing_path = parsed.path[: -len(".full.pdf")]
        candidates.append(parsed._replace(path=landing_path, query="", fragment="").geturl())
    elif lowered_path.endswith(".pdf"):
        landing_path = parsed.path[:-4]
        candidates.append(parsed._replace(path=landing_path, query="", fragment="").geturl())

    return unique_ordered(candidate for candidate in candidates if candidate and candidate != article_url)


def infer_doi_from_article_url(article_url: str) -> str | None:
    parsed = urlparse(article_url)
    path = parsed.path.rstrip("/")
    host = parsed.netloc.lower()

    if "nature.com" in host and "/articles/" in path:
        suffix = path.split("/articles/", 1)[1].strip("/")
        if suffix:
            return f"10.1038/{suffix}"

    if "springer.com" in host and "/article/" in path:
        suffix = path.split("/article/", 1)[1].strip("/")
        if suffix:
            return suffix

    return None


def extract_pmcid(url: str) -> str | None:
    parsed = urlparse(url)
    for part in parsed.path.split("/"):
        if part.upper().startswith("PMC"):
            return part.upper()
    return None


def europe_pmc_pdf_url(pmcid: str) -> str:
    pmcid = pmcid.upper()
    if not pmcid.startswith("PMC"):
        pmcid = f"PMC{pmcid}"
    return f"https://europepmc.org/articles/{pmcid}?pdf=render"


def parse_europe_pmc_search_candidates(payload: dict) -> list[str]:
    hits = payload.get("resultList", {}).get("result", [])
    candidates: list[str] = []
    for hit in hits:
        pmcid = (hit.get("pmcid") or "").upper()
        if not pmcid.startswith("PMC"):
            continue
        if (hit.get("hasPDF") or "").upper() != "Y":
            continue
        candidates.append(europe_pmc_pdf_url(pmcid))
    return unique_ordered(candidates)


def unique_ordered(values: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for value in values:
        if value not in seen:
            result.append(value)
            seen.add(value)
    return result


def cache_path_for(source: PaperSource, cache_dir: Path) -> Path:
    return cache_dir / f"{source.slug}.pdf"


def rendered_path_for(source: PaperSource, image_dir: Path) -> Path:
    return image_dir / f"{source.slug}_page1.png"


def is_pdf_file(path: Path) -> bool:
    if not path.exists() or path.stat().st_size < 5:
        return False
    with path.open("rb") as handle:
        header = handle.read(8).lstrip()
    if not header.startswith(b"%PDF"):
        return False
    try:
        with fitz.open(path) as doc:
            if doc.page_count < 1:
                return False
            first_text = doc[0].get_text()
            return not is_supplementary_pdf_text(first_text)
    except Exception:
        return False


def is_supplementary_pdf_text(text: str) -> bool:
    normalized = " ".join(text.lower().split())
    leading_excerpt = normalized[:240]
    supplementary_headings = [
        "supplementary information",
        "supplementary methods",
        "supplementary material",
        "reporting summary",
    ]
    if any(leading_excerpt.startswith(marker) for marker in supplementary_headings):
        return True
    return "in the format provided by the authors and unedited" in normalized


def looks_like_pdf_bytes(content: bytes) -> bool:
    return content[:1024].lstrip().startswith(b"%PDF")


def is_supplementary_pdf_link(url: str, label: str = "") -> bool:
    needle = f"{url} {label}".lower()
    markers = [
        "supplementary",
        "moesm",
        "_esm",
        "/esm/",
        "mediaobjects",
        "reporting summary",
        "peer review file",
    ]
    return any(marker in needle for marker in markers)


def parse_citation_doi_from_html(html: str) -> str | None:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup.find_all("meta"):
        name = (tag.get("name") or tag.get("property") or tag.get("itemprop") or "").lower()
        content = (tag.get("content") or "").strip()
        if content and name in {"citation_doi", "dc.identifier", "doi"}:
            return content
    return None


def parse_pdf_links_from_html(html: str, base_url: str) -> list[str]:
    soup = BeautifulSoup(html, "html.parser")
    found: list[str] = []

    for tag in soup.find_all("meta"):
        name = (tag.get("name") or tag.get("property") or "").lower()
        content = tag.get("content") or ""
        if "pdf" in name and ".pdf" in content.lower() and not is_supplementary_pdf_link(content):
            found.append(urljoin(base_url, content))

    for tag in soup.find_all("a", href=True):
        href = tag.get("href", "")
        text = tag.get_text(" ", strip=True).lower()
        href_lower = href.lower()
        if is_supplementary_pdf_link(href, text):
            continue
        if ".pdf" in href_lower or "/pdf" in href_lower or "download pdf" in text:
            found.append(urljoin(base_url, href))

    return unique_ordered(found)


def normalize_pmc_ftp_url(url: str) -> str:
    if not url.startswith("ftp://ftp.ncbi.nlm.nih.gov/pub/pmc/"):
        return url
    path = url[len("ftp://ftp.ncbi.nlm.nih.gov/pub/pmc/") :]
    if not path.startswith("deprecated/"):
        path = f"deprecated/{path}"
    return f"https://ftp.ncbi.nlm.nih.gov/pub/pmc/{path}"


def parse_pmc_oa_link_candidates(xml_text: str) -> list[str]:
    try:
        root = ET.fromstring(xml_text)
    except ET.ParseError:
        return []

    candidates: list[str] = []
    for link in root.findall(".//link"):
        fmt = link.attrib.get("format", "").lower()
        if fmt not in {"pdf", "tgz"}:
            continue
        href = link.attrib.get("href")
        if href:
            candidates.append(normalize_pmc_ftp_url(href))
    return unique_ordered(candidates)


def should_attempt_interactive_login(errors: list[str]) -> bool:
    blob = " ".join(errors).lower()
    login_signals = [
        "returned html instead of pdf",
        "403",
        "forbidden",
        "access through your institution",
        "sign in",
        "login",
    ]
    return any(signal in blob for signal in login_signals)


def request_url(session: requests.Session, url: str, timeout: int) -> requests.Response:
    response = session.get(
        url,
        headers={
            "User-Agent": USER_AGENT,
            "Accept": "application/pdf,text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        },
        timeout=timeout,
        allow_redirects=True,
    )
    response.raise_for_status()
    return response


def request_binary(session: requests.Session, url: str, timeout: int) -> tuple[bytes, str, str]:
    if url.startswith("ftp://"):
        url = normalize_pmc_ftp_url(url)
    if url.startswith("ftp://"):
        request = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
        with urllib.request.urlopen(request, timeout=timeout) as response:
            return response.read(), response.geturl(), response.headers.get("content-type", "")
    response = request_url(session, url, timeout)
    return response.content, response.url, response.headers.get("content-type", "")


def request_json(session: requests.Session, url: str, timeout: int) -> dict:
    response = session.get(
        url,
        headers={
            "User-Agent": USER_AGENT,
            "Accept": "application/json,text/plain,*/*",
        },
        timeout=timeout,
        allow_redirects=True,
    )
    response.raise_for_status()
    return response.json()


def pmc_oa_candidates_for_pmcid(pmcid: str, session: requests.Session, timeout: int) -> list[str]:
    if not pmcid:
        return []
    oa_url = f"https://www.ncbi.nlm.nih.gov/pmc/utils/oa/oa.fcgi?id={pmcid.upper()}"
    try:
        response = request_url(session, oa_url, timeout)
    except Exception:
        return []
    return parse_pmc_oa_link_candidates(response.text)


def pmc_oa_pdf_candidates(article_url: str, session: requests.Session, timeout: int) -> list[str]:
    parsed = urlparse(article_url)
    if "ncbi.nlm.nih.gov" not in parsed.netloc.lower():
        return []
    parts = [part for part in parsed.path.split("/") if part.upper().startswith("PMC")]
    if not parts:
        return []
    return pmc_oa_candidates_for_pmcid(parts[0], session, timeout)


def parse_europe_pmc_search_pmcids(payload: dict) -> list[str]:
    hits = payload.get("resultList", {}).get("result", [])
    pmcids: list[str] = []
    for hit in hits:
        pmcid = (hit.get("pmcid") or "").upper()
        if pmcid.startswith("PMC"):
            pmcids.append(pmcid)
    return unique_ordered(pmcids)


def europe_pmc_candidates(
    article_url: str,
    session: requests.Session,
    timeout: int,
    *,
    html_pages: list[tuple[str, str]] | None = None,
) -> list[str]:
    pmcids: list[str] = []
    pmcid = extract_pmcid(article_url)
    if pmcid:
        pmcids.append(pmcid)

    doi_candidates = [infer_doi_from_article_url(article_url)]
    for html, _base_url in html_pages or []:
        doi_candidates.append(parse_citation_doi_from_html(html))

    candidates: list[str] = []
    for doi in unique_ordered(value for value in doi_candidates if value):
        query = quote(f"DOI:{doi}", safe=":/")
        api_url = (
            "https://www.ebi.ac.uk/europepmc/webservices/rest/search"
            f"?query={query}&format=json&pageSize=5"
        )
        try:
            payload = request_json(session, api_url, timeout)
        except Exception:
            continue
        candidates.extend(parse_europe_pmc_search_candidates(payload))
        pmcids.extend(parse_europe_pmc_search_pmcids(payload))

    for pmcid in unique_ordered(pmcids):
        candidates.extend(pmc_oa_candidates_for_pmcid(pmcid, session, timeout))

    return unique_ordered(candidates)


def looks_like_tar_gz(content: bytes, content_type: str, url: str) -> bool:
    lowered_url = url.lower()
    lowered_type = content_type.lower()
    return (
        content.startswith(b"\x1f\x8b")
        and (lowered_url.endswith(".tar.gz") or lowered_url.endswith(".tgz") or "gzip" in lowered_type)
    )


def extract_pdf_bytes_from_pmc_archive(content: bytes) -> bytes | None:
    try:
        with tarfile.open(fileobj=io.BytesIO(content), mode="r:gz") as archive:
            for member in archive.getmembers():
                if not member.isfile() or not member.name.lower().endswith(".pdf"):
                    continue
                extracted = archive.extractfile(member)
                if extracted is None:
                    continue
                pdf_bytes = extracted.read()
                if looks_like_pdf_bytes(pdf_bytes):
                    return pdf_bytes
    except tarfile.TarError:
        return None
    return None


def extract_pdf_candidate_bytes(content: bytes, content_type: str, final_url: str) -> bytes | None:
    if looks_like_pdf_bytes(content):
        return content
    if looks_like_tar_gz(content, content_type, final_url):
        return extract_pdf_bytes_from_pmc_archive(content)
    return None


def save_valid_pdf_bytes(destination: Path, content: bytes) -> bool:
    destination.write_bytes(content)
    if is_pdf_file(destination):
        return True
    try:
        destination.unlink()
    except FileNotFoundError:
        pass
    return False


def try_download_candidates(
    candidate_urls: list[str],
    destination: Path,
    session: requests.Session,
    timeout: int,
    success_status: str,
) -> tuple[dict | None, list[str], list[tuple[str, str]]]:
    errors: list[str] = []
    html_pages: list[tuple[str, str]] = []

    for candidate in candidate_urls:
        try:
            content, final_url, content_type = request_binary(session, candidate, timeout)
            pdf_bytes = extract_pdf_candidate_bytes(content, content_type, final_url)
            if pdf_bytes is not None:
                if save_valid_pdf_bytes(destination, pdf_bytes):
                    return (
                        {
                            "status": success_status,
                            "path": str(destination),
                            "url": final_url,
                            "bytes": len(pdf_bytes),
                        },
                        errors,
                        html_pages,
                    )
                errors.append(f"{candidate}: PDF appears to be supplementary material")
                continue
            if "html" in content_type.lower() or b"<html" in content[:2048].lower():
                html_pages.append((content.decode("utf-8", errors="replace"), final_url))
                errors.append(f"{candidate}: returned HTML instead of PDF")
            else:
                errors.append(f"{candidate}: not a PDF, content-type={content_type}")
        except Exception as exc:
            errors.append(f"{candidate}: {type(exc).__name__}: {exc}")

    return None, errors, html_pages


def requests_session_from_browser_cookies(cookies: list[dict]) -> requests.Session:
    session = requests.Session()
    for cookie in cookies:
        name = cookie.get("name")
        value = cookie.get("value")
        domain = cookie.get("domain")
        path = cookie.get("path") or "/"
        if not name or value is None or not domain:
            continue
        expires = cookie.get("expires")
        if isinstance(expires, float):
            expires = int(expires)
        session.cookies.set(
            name,
            value,
            domain=domain,
            path=path,
            secure=bool(cookie.get("secure", False)),
            expires=expires if isinstance(expires, int) and expires > 0 else None,
        )
    return session


def import_playwright_sync():
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as exc:
        raise RuntimeError(
            "Interactive login requires Playwright. Install it with "
            "`python -m pip install playwright` and then run `playwright install chromium`."
        ) from exc
    return sync_playwright


def download_pdf_via_interactive_login(
    source: PaperSource,
    destination: Path,
    candidate_urls: list[str],
    timeout: int,
    config: InteractiveLoginConfig,
) -> tuple[dict | None, list[str]]:
    sync_playwright = import_playwright_sync()
    config.profile_dir.mkdir(parents=True, exist_ok=True)
    errors: list[str] = []

    print(
        f"[login] Opening browser for {source.slug}. "
        f"Please complete institution login; the script will retry automatically.",
        flush=True,
    )

    with sync_playwright() as playwright:
        try:
            context = playwright.chromium.launch_persistent_context(
                str(config.profile_dir.resolve()),
                headless=False,
                accept_downloads=True,
                viewport={"width": 1440, "height": 960},
            )
        except Exception as exc:
            raise RuntimeError(
                "Could not start Playwright Chromium. Run `playwright install chromium` and try again."
            ) from exc

        try:
            login_page = context.pages[0] if context.pages else context.new_page()
            try:
                login_page.goto(source.article_url, wait_until="domcontentloaded", timeout=60000)
            except Exception as exc:
                errors.append(f"{source.article_url}: browser open failed: {type(exc).__name__}: {exc}")

            deadline = time.time() + config.timeout_seconds
            next_notice_at = 0.0
            while time.time() < deadline:
                auth_session = requests_session_from_browser_cookies(context.cookies())
                result, retry_errors, _html_pages = try_download_candidates(
                    candidate_urls=candidate_urls,
                    destination=destination,
                    session=auth_session,
                    timeout=timeout,
                    success_status="downloaded_after_login",
                )
                if result is not None:
                    return result, errors

                errors = retry_errors
                now = time.time()
                if now >= next_notice_at:
                    remaining = max(0, int(deadline - now))
                    print(
                        f"[login] Waiting for completed login on {source.slug}; "
                        f"retrying in {config.poll_interval_seconds}s ({remaining}s left).",
                        flush=True,
                    )
                    next_notice_at = now + 30
                time.sleep(config.poll_interval_seconds)
        finally:
            context.close()

    errors.append(f"Interactive login timed out after {config.timeout_seconds} seconds")
    return None, errors


def download_pdf(
    source: PaperSource,
    cache_dir: Path,
    session: requests.Session,
    timeout: int = 45,
    interactive_login: InteractiveLoginConfig | None = None,
) -> dict:
    cache_dir.mkdir(parents=True, exist_ok=True)
    destination = cache_path_for(source, cache_dir)
    if is_pdf_file(destination):
        return {
            "slug": source.slug,
            "status": "cached",
            "path": str(destination),
            "url": None,
            "bytes": destination.stat().st_size,
        }

    candidates = unique_ordered(
        [candidate for candidate in [source.pdf_url, *infer_pdf_candidates(source.article_url)] if candidate]
    )
    errors: list[str] = []
    html_pages: list[tuple[str, str]] = []

    first_result, first_errors, first_html_pages = try_download_candidates(
        candidate_urls=candidates,
        destination=destination,
        session=session,
        timeout=timeout,
        success_status="downloaded",
    )
    if first_result is not None:
        return {"slug": source.slug, **first_result}
    errors.extend(first_errors)
    html_pages.extend(first_html_pages)

    for landing_url in infer_landing_page_candidates(source.article_url):
        try:
            response = request_url(session, landing_url, timeout)
            if "html" in response.headers.get("content-type", "").lower():
                html_pages.append((response.text, response.url))
        except Exception as exc:
            errors.append(f"{landing_url}: landing page probe failed: {type(exc).__name__}: {exc}")

    extra_candidates: list[str] = []
    for html, base_url in html_pages:
        extra_candidates.extend(parse_pdf_links_from_html(html, base_url))
    if not html_pages:
        try:
            response = request_url(session, source.article_url, timeout)
            extra_candidates.extend(parse_pdf_links_from_html(response.text, response.url))
        except Exception as exc:
            errors.append(f"{source.article_url}: html parse failed: {type(exc).__name__}: {exc}")
    extra_candidates.extend(pmc_oa_pdf_candidates(source.article_url, session, timeout))
    extra_candidates.extend(europe_pmc_candidates(source.article_url, session, timeout, html_pages=html_pages))

    browser_candidates = unique_ordered([*candidates, *extra_candidates])
    secondary_candidates = [candidate for candidate in unique_ordered(extra_candidates) if candidate not in candidates]
    second_result, second_errors, _second_html_pages = try_download_candidates(
        candidate_urls=secondary_candidates,
        destination=destination,
        session=session,
        timeout=timeout,
        success_status="downloaded_from_html",
    )
    if second_result is not None:
        return {"slug": source.slug, **second_result}
    errors.extend(second_errors)

    if interactive_login and interactive_login.enabled and should_attempt_interactive_login(errors):
        login_result, login_errors = download_pdf_via_interactive_login(
            source=source,
            destination=destination,
            candidate_urls=browser_candidates,
            timeout=timeout,
            config=interactive_login,
        )
        if login_result is not None:
            return {"slug": source.slug, **login_result}
        errors.extend(login_errors)

    return {
        "slug": source.slug,
        "status": "failed",
        "path": str(destination),
        "url": None,
        "errors": errors,
        "manual_save_as": str(destination),
    }


def validate_all_sources_available(
    sources: list[PaperSource], cache_dir: Path, allow_missing: bool
) -> list[PaperSource]:
    missing = [source for source in sources if not is_pdf_file(cache_path_for(source, cache_dir))]
    if missing and not allow_missing:
        return missing
    return missing


def render_first_page(source: PaperSource, pdf_path: Path, image_dir: Path, zoom: float = 2.4) -> Path:
    image_dir.mkdir(parents=True, exist_ok=True)
    output = rendered_path_for(source, image_dir)
    if output.exists() and output.stat().st_size > 0:
        return output

    with fitz.open(pdf_path) as doc:
        if doc.page_count < 1:
            raise ValueError(f"{pdf_path} has no pages")
        page = doc.load_page(0)
        matrix = fitz.Matrix(zoom, zoom)
        pixmap = page.get_pixmap(matrix=matrix, alpha=False)
        pixmap.save(output)
    return output


def fit_to_page(image: Image.Image, page_width: int, page_height: int) -> Image.Image:
    source = image.convert("RGB")
    canvas = Image.new("RGB", (page_width, page_height), "white")
    scale = min(page_width / source.width, page_height / source.height)
    new_size = (max(1, int(source.width * scale)), max(1, int(source.height * scale)))
    resized = source.resize(new_size, Image.Resampling.LANCZOS)
    x = (page_width - new_size[0]) // 2
    y = (page_height - new_size[1]) // 2
    canvas.paste(resized, (x, y))
    return canvas


def add_shadow_and_border(page: Image.Image, shadow: int = 18, border: int = 2) -> Image.Image:
    page_rgba = page.convert("RGBA")
    width, height = page_rgba.size
    output = Image.new("RGBA", (width + shadow * 2, height + shadow * 2), (255, 255, 255, 0))

    shadow_layer = Image.new("RGBA", (width, height), (0, 0, 0, 70))
    shadow_layer = shadow_layer.filter(ImageFilter.GaussianBlur(radius=10))
    output.paste(shadow_layer, (shadow + 9, shadow + 10), shadow_layer)
    output.paste(page_rgba, (shadow, shadow), page_rgba)

    draw = ImageDraw.Draw(output)
    draw.rectangle(
        [shadow, shadow, shadow + width - 1, shadow + height - 1],
        outline=(45, 45, 45, 255),
        width=border,
    )
    return output


def build_draw_plan(
    item_count: int,
    page_width: int,
    page_height: int,
    left_front: tuple[int, int],
    right_front: tuple[int, int],
    gap_x: int,
    gap_y: int,
) -> list[DrawPlacement]:
    if item_count == 0:
        return []
    group_size = 6 if item_count >= 12 else max(1, (item_count + 1) // 2)
    placements: list[DrawPlacement] = []
    group_fronts = [left_front, right_front]

    for group_index, start in enumerate(range(0, item_count, group_size)):
        group_items = list(range(start, min(start + group_size, item_count)))
        front_x, front_y = group_fronts[min(group_index, 1)]
        last_layer = len(group_items) - 1
        for layer_index, source_index in enumerate(group_items):
            offset_layers = last_layer - layer_index
            x = front_x + offset_layers * gap_x
            y = front_y - offset_layers * gap_y
            placements.append(
                DrawPlacement(
                    source_index=source_index,
                    group_index=group_index,
                    layer_index=layer_index,
                    paste_xy=(x, y),
                )
            )
    return placements


def make_collage(
    sources: list[PaperSource],
    image_dir: Path,
    output_path: Path,
    pdf_output_path: Path | None,
    pptx_output_path: Path | None,
    canvas_width: int,
    canvas_height: int,
    page_width: int,
    page_height: int,
    gap_x: int,
    gap_y: int,
    background: str,
) -> dict:
    canvas = Image.new("RGB", (canvas_width, canvas_height), background)
    rendered_pages = build_rendered_pages(
        sources=sources,
        image_dir=image_dir,
        page_width=page_width,
        page_height=page_height,
        gap_x=gap_x,
        gap_y=gap_y,
    )

    layer_report = []
    for rendered_page in rendered_pages:
        with Image.open(rendered_page.framed_path) as framed:
            canvas.paste(framed.convert("RGB"), rendered_page.paste_xy, framed)
        layer_report.append(
            {
                **asdict(rendered_page.placement),
                "slug": rendered_page.source.slug,
                "framed_page": str(rendered_page.framed_path),
                "paste_xy": list(rendered_page.paste_xy),
                "framed_size": list(rendered_page.size),
            }
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(output_path)
    if pdf_output_path:
        pdf_output_path.parent.mkdir(parents=True, exist_ok=True)
        canvas.save(pdf_output_path, "PDF", resolution=150.0)
    if pptx_output_path:
        make_editable_pptx(
            rendered_pages=rendered_pages,
            pptx_output_path=pptx_output_path,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            background=background,
        )

    return {
        "output": str(output_path),
        "pdf_output": str(pdf_output_path) if pdf_output_path else None,
        "pptx_output": str(pptx_output_path) if pptx_output_path else None,
        "canvas": [canvas_width, canvas_height],
        "layers": layer_report,
    }


def build_rendered_pages(
    sources: list[PaperSource],
    image_dir: Path,
    page_width: int,
    page_height: int,
    gap_x: int,
    gap_y: int,
) -> list[RenderedPage]:
    placements = build_draw_plan(
        item_count=len(sources),
        page_width=page_width,
        page_height=page_height,
        left_front=(150, 365),
        right_front=(955, 365),
        gap_x=gap_x,
        gap_y=gap_y,
    )
    framed_dir = image_dir / "_framed_pages"
    framed_dir.mkdir(parents=True, exist_ok=True)

    rendered_pages: list[RenderedPage] = []
    for placement in placements:
        source = sources[placement.source_index]
        rendered_path = rendered_path_for(source, image_dir)
        with Image.open(rendered_path) as img:
            fitted = fit_to_page(img, page_width, page_height)
        framed = add_shadow_and_border(fitted)
        framed_path = framed_dir / f"{source.slug}_framed.png"
        framed.save(framed_path)
        paste_x, paste_y = placement.paste_xy
        paste_x -= 18
        paste_y -= 18
        rendered_pages.append(
            RenderedPage(
                source=source,
                placement=placement,
                framed_path=framed_path,
                paste_xy=(paste_x, paste_y),
                size=framed.size,
            )
        )
    return rendered_pages


def px_to_emu(px: int) -> int:
    return int(px * 9525)


def background_to_rgb(background: str) -> tuple[int, int, int]:
    return ImageColor.getrgb(background)


def make_editable_pptx(
    rendered_pages: list[RenderedPage],
    pptx_output_path: Path,
    canvas_width: int,
    canvas_height: int,
    background: str,
) -> str:
    pptx_output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = px_to_emu(canvas_width)
    presentation.slide_height = px_to_emu(canvas_height)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    fill = slide.background.fill
    fill.solid()
    red, green, blue = background_to_rgb(background)
    fill.fore_color.rgb = RGBColor(red, green, blue)

    for rendered_page in rendered_pages:
        left, top = rendered_page.paste_xy
        width, height = rendered_page.size
        picture = slide.shapes.add_picture(
            str(rendered_page.framed_path),
            px_to_emu(left),
            px_to_emu(top),
            width=px_to_emu(width),
            height=px_to_emu(height),
        )
        picture.name = rendered_page.source.slug

    presentation.save(str(pptx_output_path))
    return str(pptx_output_path)


def write_json(path: Path, data: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")


def write_missing_report(path: Path, missing: list[PaperSource], cache_dir: Path) -> None:
    lines = [
        "Missing real PDFs. Save each file to the listed path, then rerun the script.",
        "",
    ]
    for source in missing:
        lines.extend(
            [
                f"- {source.slug} {source.refs}",
                f"  Article URL: {source.article_url}",
                f"  Candidate PDF URLs: {', '.join([candidate for candidate in [source.pdf_url, *infer_pdf_candidates(source.article_url)] if candidate])}",
                f"  Save as: {(cache_dir / (source.slug + '.pdf')).resolve()}",
            ]
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(lines), encoding="utf-8")


def run_collage_pipeline(
    sources: list[PaperSource],
    *,
    cache_dir: Path,
    image_dir: Path,
    output_path: Path,
    pdf_output_path: Path | None,
    pptx_output_path: Path | None,
    canvas_width: int,
    canvas_height: int,
    page_width: int,
    page_height: int,
    gap_x: int,
    gap_y: int,
    background: str,
    allow_missing: bool,
    skip_download: bool,
    timeout: int,
    interactive_login: InteractiveLoginConfig,
) -> int:
    report_dir = output_path.parent

    download_report: list[dict] = []
    session = requests.Session()
    if not skip_download:
        for source in sources:
            print(f"[download] {source.slug} {source.article_url}", flush=True)
            result = download_pdf(
                source,
                cache_dir,
                session,
                timeout=timeout,
                interactive_login=interactive_login,
            )
            download_report.append(result)
            print(f"  -> {result['status']}", flush=True)
            time.sleep(0.2)
    else:
        download_report.append({"status": "skipped", "reason": "--skip-download"})

    write_json(report_dir / "download_report.json", download_report)

    missing = validate_all_sources_available(sources, cache_dir, allow_missing=allow_missing)
    if missing and not allow_missing:
        write_missing_report(report_dir / "missing_pdfs.txt", missing, cache_dir)
        print(f"Missing {len(missing)} PDFs. See {report_dir / 'missing_pdfs.txt'}", file=sys.stderr)
        raise MissingPdfError(missing)

    usable_sources = [source for source in sources if is_pdf_file(cache_path_for(source, cache_dir))]
    if not usable_sources:
        print("No valid PDF files are available. Nothing was generated.", file=sys.stderr)
        return 2

    render_report: list[dict] = []
    for source in usable_sources:
        pdf_path = cache_path_for(source, cache_dir)
        print(f"[render] {source.slug}", flush=True)
        rendered_path = render_first_page(source, pdf_path, image_dir)
        render_report.append(
            {
                "slug": source.slug,
                "pdf": str(pdf_path),
                "page1_image": str(rendered_path),
                "status": "rendered",
            }
        )
    write_json(report_dir / "render_report.json", render_report)

    collage_report = make_collage(
        sources=usable_sources,
        image_dir=image_dir,
        output_path=output_path,
        pdf_output_path=pdf_output_path,
        pptx_output_path=pptx_output_path,
        canvas_width=canvas_width,
        canvas_height=canvas_height,
        page_width=page_width,
        page_height=page_height,
        gap_x=gap_x,
        gap_y=gap_y,
        background=background,
    )
    write_json(report_dir / "collage_report.json", collage_report)
    print(f"Generated {output_path}")
    if pdf_output_path:
        print(f"Generated {pdf_output_path}")
    if pptx_output_path:
        print(f"Generated {pptx_output_path}")
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Create a collage from real paper PDF first pages.")
    parser.add_argument("--output", default="output/real_paper_cover_collage_16x9.png")
    parser.add_argument("--pdf-output", default=None)
    parser.add_argument("--pptx-output", default=None)
    parser.add_argument("--cache-dir", default="pdf_cache")
    parser.add_argument("--image-dir", default="page1_cache")
    parser.add_argument("--canvas-width", type=int, default=1920)
    parser.add_argument("--canvas-height", type=int, default=1080)
    parser.add_argument("--page-width", type=int, default=390)
    parser.add_argument("--page-height", type=int, default=520)
    parser.add_argument("--gap-x", type=int, default=70)
    parser.add_argument("--gap-y", type=int, default=52)
    parser.add_argument("--background", default="white")
    parser.add_argument("--allow-missing", action="store_true")
    parser.add_argument("--skip-download", action="store_true")
    parser.add_argument("--timeout", type=int, default=45)
    parser.add_argument("--interactive-login", action="store_true")
    parser.add_argument("--login-profile-dir", default=".playwright_login_profile")
    parser.add_argument("--login-timeout", type=int, default=900)
    parser.add_argument("--login-poll-interval", type=int, default=10)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    cache_dir = Path(args.cache_dir)
    image_dir = Path(args.image_dir)
    output_path = Path(args.output)
    pdf_output_path = Path(args.pdf_output) if args.pdf_output else output_path.with_suffix(".pdf")
    pptx_output_path = Path(args.pptx_output) if args.pptx_output else None
    interactive_login = InteractiveLoginConfig(
        enabled=args.interactive_login,
        profile_dir=Path(args.login_profile_dir),
        timeout_seconds=args.login_timeout,
        poll_interval_seconds=args.login_poll_interval,
    )
    return run_collage_pipeline(
        PDF_SOURCES,
        cache_dir=cache_dir,
        image_dir=image_dir,
        output_path=output_path,
        pdf_output_path=pdf_output_path,
        pptx_output_path=pptx_output_path,
        canvas_width=args.canvas_width,
        canvas_height=args.canvas_height,
        page_width=args.page_width,
        page_height=args.page_height,
        gap_x=args.gap_x,
        gap_y=args.gap_y,
        background=args.background,
        allow_missing=args.allow_missing,
        skip_download=args.skip_download,
        timeout=args.timeout,
        interactive_login=interactive_login,
    )


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except MissingPdfError as exc:
        for source in exc.missing:
            print(f"- {source.slug}: {source.article_url}", file=sys.stderr)
        raise SystemExit(1)
